import io
import time
import urllib.parse
import streamlit as st
from docx import Document
from pptx import Presentation

# ==========================================
# SAFE IMPORTS & FALLBACKS
# ==========================================
try:
    from gtts import gTTS
    HAS_GTTS = True
except ImportError:
    HAS_GTTS = False

# ==========================================
# 1. PAGE CONFIGURATION & MODERN STYLING
# ==========================================
st.set_page_config(
    page_title="Ultimate Exam Prep & Goal Hub",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom Glassmorphic & Modern CSS Styling
st.markdown("""
    <style>
    /* Dark Theme Background & Fonts */
    .stApp {
        background: linear-gradient(135deg, #0f172a 0%, #1e1b4b 100%);
        color: #f8fafc;
    }
    
    /* Custom Card Styling */
    .glass-card {
        background: rgba(30, 41, 59, 0.7);
        backdrop-filter: blur(10px);
        border-radius: 16px;
        padding: 24px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        margin-bottom: 20px;
        box-shadow: 0 10px 25px -5px rgba(0, 0, 0, 0.3);
    }
    
    .badge {
        display: inline-block;
        padding: 4px 12px;
        border-radius: 9999px;
        font-size: 12px;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 0.05em;
        margin-right: 8px;
    }
    .badge-primary { background-color: #3b82f6; color: white; }
    .badge-success { background-color: #10b981; color: white; }
    .badge-purple { background-color: #8b5cf6; color: white; }

    /* Custom Buttons */
    .yt-btn {
        display: inline-flex;
        align-items: center;
        justify-content: center;
        background: linear-gradient(135deg, #ef4444 0%, #dc2626 100%);
        color: white !important;
        padding: 10px 20px;
        border-radius: 10px;
        text-decoration: none;
        font-weight: 600;
        margin-top: 12px;
        box-shadow: 0 4px 12px rgba(239, 68, 68, 0.3);
        transition: all 0.3s ease;
    }
    .yt-btn:hover { transform: translateY(-2px); box-shadow: 0 6px 16px rgba(239, 68, 68, 0.4); text-decoration: none; }
    
    .official-btn {
        display: inline-block;
        background: linear-gradient(135deg, #10b981 0%, #059669 100%);
        color: white !important;
        padding: 12px 24px;
        border-radius: 10px;
        text-decoration: none;
        font-weight: 700;
        margin-bottom: 20px;
        box-shadow: 0 4px 12px rgba(16, 185, 129, 0.3);
    }
    .official-btn:hover { text-decoration: none; transform: translateY(-2px); }

    .download-link-btn {
        display: inline-block;
        background-color: #2563eb;
        color: white !important;
        padding: 8px 16px;
        border-radius: 8px;
        text-decoration: none;
        font-weight: 600;
        font-size: 13px;
    }

    .share-btn-wa {
        display: inline-flex;
        align-items: center;
        justify-content: center;
        width: 100%;
        background-color: #25D366;
        color: white !important;
        padding: 12px;
        border-radius: 10px;
        text-decoration: none;
        font-weight: bold;
        text-align: center;
    }
    .share-btn-email {
        display: inline-flex;
        align-items: center;
        justify-content: center;
        width: 100%;
        background-color: #ea4335;
        color: white !important;
        padding: 12px;
        border-radius: 10px;
        text-decoration: none;
        font-weight: bold;
        text-align: center;
    }
    
    .roadmap-step {
        border-left: 4px solid #3b82f6;
        background: rgba(30, 41, 59, 0.4);
        padding: 16px;
        border-radius: 0 12px 12px 0;
        margin-bottom: 16px;
    }
    </style>
""", unsafe_allow_html=True)


# ==========================================
# 2. DOCX, PPTX & AUDIO EXPORT GENERATORS
# ==========================================
def create_docx_plan(exam_name, age_group, daily_hours, prep_months, plan_text):
    doc = Document()
    doc.add_heading(f"Customized Goal Plan: {exam_name}", level=0)
    
    p = doc.add_paragraph()
    p.add_run("Target Exam: ").bold = True
    p.add_run(f"{exam_name}\n")
    p.add_run("Aspirant Profile / Age Group: ").bold = True
    p.add_run(f"{age_group}\n")
    p.add_run("Timeline: ").bold = True
    p.add_run(f"{prep_months} Months\n")
    p.add_run("Daily Capacity: ").bold = True
    p.add_run(f"{daily_hours} Hours/day")

    doc.add_heading("Detailed Strategy & Timetable", level=1)
    
    for line in plan_text.split('\n'):
        if line.strip():
            doc.add_paragraph(line.strip())
            
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

def create_pptx_plan(exam_name, age_group, daily_hours, prep_months, plan_text):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Goal Plan: {exam_name}"
    subtitle.text = f"Profile: {age_group}\nTimeline: {prep_months} Months | Daily: {daily_hours} Hrs"
    
    # Plan Content Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Execution Strategy & Daily Routine"
    
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    lines = [l.strip() for l in plan_text.split('\n') if l.strip()]
    for i, line in enumerate(lines[:10]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

def text_to_speech(text):
    """Generates audio bytes from text using gTTS if installed."""
    if not HAS_GTTS:
        return None
    tts = gTTS(text=text, lang='en')
    fp = io.BytesIO()
    tts.write_to_fp(fp)
    fp.seek(0)
    return fp


# ==========================================
# 3. EXPANDED EXAM & MOCK PRACTICE DATABASE
# ==========================================
EXAMS_DATA = {
    "SSC CGL": {
        "full_name": "SSC Combined Graduate Level",
        "icon": "🏛️",
        "category": "SSC Group B & C",
        "description": "Premier national recruitment exam for Inspector, Assistant Section Officer, and Tax Officer posts.",
        "official_pyq_portal": "https://ssc.gov.in/",
        "roadmap": [
            "**Phase 1: Concepts & Core Math (Months 1-3):** Cover Advance & Arithmetic Maths, English Grammar rules, and Reasoning topics.",
            "**Phase 2: Speed Building & Speed Math (Months 4-5):** Solve 5,000+ Previous Year Questions (PYQs) with speed tricks.",
            "**Phase 3: Tier-1 & Tier-2 Mock Drills (Month 6):** Daily sectional tests, full mocks, and typing practice."
        ],
        "top_teachers": [
            {"subject": "Quantitative Aptitude (Maths)", "teacher": "Gagan Pratap Sir", "channel": "Gagan Pratap Maths", "video_title": "Complete Maths Marathon", "youtube_link": "https://www.youtube.com/results?search_query=ssc+cgl+maths+gagan+pratap+marathon"},
            {"subject": "Maths Practice Drills", "teacher": "Aditya Ranjan Sir", "channel": "Rankers Gurukul", "video_title": "CGL Practice Series", "youtube_link": "https://www.youtube.com/results?search_query=aditya+ranjan+sir+cgl+maths"},
            {"subject": "English Language", "teacher": "Nimisha Bansal Ma'am", "channel": "Nimisha Bansal", "video_title": "120 Rules of Grammar Complete", "youtube_link": "https://www.youtube.com/results?search_query=120+rules+of+grammar+nimisha+bansal"},
            {"subject": "General Studies & Static GK", "teacher": "Parmar SSC", "channel": "Parmar SSC", "video_title": "Complete Static GK Batch", "youtube_link": "https://www.youtube.com/results?search_query=parmar+ssc+gk"}
        ],
        "pyqs": [
            {"year": "2024", "title": "SSC CGL Tier-1 All Shifts Question Papers", "link": "https://ssc.gov.in/", "format": "Official PDF Archive"},
            {"year": "2023", "title": "SSC CGL Tier-2 (Mains) Official Question Paper", "link": "https://ssc.gov.in/", "format": "Official PDF Archive"}
        ],
        "mocks": [
            {"question": "Which article of the Indian Constitution empowers the President to declare a National Emergency?", "options": ["Article 352", "Article 356", "Article 360", "Article 370"], "answer": "Article 352", "explanation": "Article 352 deals with National Emergency due to war, external aggression, or armed rebellion."},
            {"question": "If the length of a rectangle is increased by 20% and breadth is decreased by 10%, what is the net change in area?", "options": ["8% increase", "10% increase", "5% decrease", "12% increase"], "answer": "8% increase", "explanation": "Net change = +20 - 10 + ((20 * -10)/100) = 10 - 2 = +8% increase."},
            {"question": "Find the correctly spelt word:", "options": ["Accomodation", "Accommodation", "Acommodation", "Accommodatoin"], "answer": "Accommodation", "explanation": "The correct spelling is 'Accommodation' with double 'c' and double 'm'."},
            {"question": "Select the option that is related to the third word in the same way: Book : Publisher :: Film : ?", "options": ["Director", "Producer", "Actor", "Writer"], "answer": "Producer", "explanation": "A publisher finances/produces a book, similarly a producer finances/produces a film."},
            {"question": "Who was the first Law Minister of Independent India?", "options": ["Dr. B.R. Ambedkar", "Jawaharlal Nehru", "Sardar Patel", "Maulana Azad"], "answer": "Dr. B.R. Ambedkar", "explanation": "Dr. B.R. Ambedkar served as the first Law and Justice Minister of India."},
            {"question": "Which gland in the human body is called the 'Master Gland'?", "options": ["Pituitary Gland", "Thyroid Gland", "Adrenal Gland", "Pancreas"], "answer": "Pituitary Gland", "explanation": "The Pituitary Gland regulates the operations of many other endocrine glands."}
        ]
    },
    "SSC CHSL": {
        "full_name": "SSC Combined Higher Secondary Level",
        "icon": "📜",
        "category": "SSC 10+2 Level",
        "description": "National exam for Lower Division Clerk (LDC), JSA, and Data Entry Operators (DEO).",
        "official_pyq_portal": "https://ssc.gov.in/",
        "roadmap": [
            "**Phase 1: Syllabus Fundamentals (Months 1-2):** Master Class 10th level Quantitative Aptitude, Basic English, and Logical Reasoning.",
            "**Phase 2: PYQ Drill & Speed Practice (Months 3-4):** Solve CHSL papers from the last 5 years.",
            "**Phase 3: Mocks & Typing Speed (Month 5):** Conduct full-length tests and daily typing practice."
        ],
        "top_teachers": [
            {"subject": "Reasoning Ability", "teacher": "Vikramjeet Sir", "channel": "Rankers Gurukul", "video_title": "SSC CHSL Reasoning Complete Playlist", "youtube_link": "https://www.youtube.com/results?search_query=ssc+chsl+reasoning+vikramjeet+sir"}
        ],
        "pyqs": [
            {"year": "2024", "title": "SSC CHSL Tier-1 Shifts Solved Papers", "link": "https://ssc.gov.in/", "format": "PDF"}
        ],
        "mocks": [
            {"question": "What is the capital of Dadra and Nagar Haveli and Daman and Diu?", "options": ["Daman", "Silvassa", "Kavaratti", "Port Blair"], "answer": "Daman", "explanation": "Daman was declared the capital of the merged UT in 2020."},
            {"question": "Simple interest on $1,000 for 2 years at 5% per annum is:", "options": ["$100", "$50", "$150", "$200"], "answer": "$100", "explanation": "SI = (P * R * T)/100 = (1000 * 5 * 2)/100 = $100."},
            {"question": "Which gas is used in fire extinguishers?", "options": ["Carbon Dioxide", "Oxygen", "Nitrogen", "Hydrogen"], "answer": "Carbon Dioxide", "explanation": "CO2 displaces oxygen around fire to extinguish it."}
        ]
    },
    "SSC GD": {
        "full_name": "SSC General Duty Constable",
        "icon": "🪖",
        "category": "SSC Defense & Paramilitary",
        "description": "Entrance exam for Constables in BSF, CISF, CRPF, SSB, ITBP, AR, and SSF.",
        "official_pyq_portal": "https://ssc.gov.in/",
        "roadmap": [
            "**Phase 1: Elementary Maths & Language (Months 1-2):** Basic arithmetic, Hindi/English grammar fundamentals.",
            "**Phase 2: GK/GS & Reasoning (Months 3-4):** Static GK, Current Affairs, and Non-Verbal Reasoning."
        ],
        "top_teachers": [
            {"subject": "General Studies & Hindi", "teacher": "Naveen Sir / Ankit Bhati Sir", "channel": "Rojgar with Ankit", "video_title": "SSC GD Complete Hindi & GS", "youtube_link": "https://www.youtube.com/results?search_query=ssc+gd+rojgar+with+ankit"}
        ],
        "pyqs": [
            {"year": "2024", "title": "SSC GD Constable Official Shifts Paper", "link": "https://ssc.gov.in/", "format": "PDF Archive"}
        ],
        "mocks": [
            {"question": "Which dance form originates from the state of Assam?", "options": ["Bihu", "Kathak", "Garba", "Ghoomar"], "answer": "Bihu", "explanation": "Bihu is the folk dance of Assam associated with the Bihu festival."},
            {"question": "Complete the series: 2, 4, 8, 16, ?", "options": ["32", "24", "64", "20"], "answer": "32", "explanation": "Each number is multiplied by 2. 16 * 2 = 32."}
        ]
    },
    "SSC JE": {
        "full_name": "SSC Junior Engineer",
        "icon": "🏗️",
        "category": "Technical Engineering",
        "description": "Recruitment exam for Junior Engineers in CPWD, MES, and Central Water Commission.",
        "official_pyq_portal": "https://ssc.gov.in/",
        "roadmap": [
            "**Phase 1: Technical Subjects (Months 1-4):** Core concepts in Civil, Electrical, or Mechanical engineering.",
            "**Phase 2: Non-Tech Aptitude (Months 5-6):** General Intelligence, Reasoning, and General Awareness."
        ],
        "top_teachers": [
            {"subject": "Civil Engineering", "teacher": "Engineers Academy", "channel": "Engineers Academy", "video_title": "SSC JE Civil Complete Revision", "youtube_link": "https://www.youtube.com/results?search_query=ssc+je+civil+engineers+academy"}
        ],
        "pyqs": [
            {"year": "2024", "title": "SSC JE Paper-1 (CBT) Official Question Paper", "link": "https://ssc.gov.in/", "format": "PDF"}
        ],
        "mocks": [
            {"question": "What is the SI unit of pressure?", "options": ["Pascal", "Joule", "Newton", "Watt"], "answer": "Pascal", "explanation": "Pascal (Pa) equals one newton per square meter."}
        ]
    },
    "SSC MTS": {
        "full_name": "SSC Multi-Tasking Staff & Havaldar",
        "icon": "📑",
        "category": "SSC 10th Level",
        "description": "Selection for non-technical General Central Service Group 'C' Posts.",
        "official_pyq_portal": "https://ssc.gov.in/",
        "roadmap": [
            "**Phase 1: Session 1 Prep (Months 1-2):** Basic Numerical Ability and Reasoning.",
            "**Phase 2: Session 2 Scoring Subjects (Months 3-4):** General Awareness and English vocabulary."
        ],
        "top_teachers": [
            {"subject": "General Knowledge & English", "teacher": "MTS Special Team", "channel": "SSC Adda247", "video_title": "SSC MTS GK & English Revision", "youtube_link": "https://www.youtube.com/results?search_query=ssc+mts+gk+english+preparation"}
        ],
        "pyqs": [
            {"year": "2024", "title": "SSC MTS All Shift Question Papers", "link": "https://ssc.gov.in/", "format": "PDF Archive"}
        ],
        "mocks": [
            {"question": "Who was the first Governor-General of Independent India?", "options": ["Lord Mountbatten", "C. Rajagopalachari", "Dr. Rajendra Prasad", "Jawaharlal Nehru"], "answer": "Lord Mountbatten", "explanation": "Lord Mountbatten was the first Governor-General of independent India."}
        ]
    },
    "UPSC CSE": {
        "full_name": "UPSC Civil Services Examination",
        "icon": "⚖️",
        "category": "Civil Services",
        "description": "India's premier exam for IAS, IPS, IFS, and Central Services.",
        "official_pyq_portal": "https://www.upsc.gov.in/examinations/previous-question-papers",
        "roadmap": [
            "**Phase 1: NCERT Foundations (Months 1-3):** Class 6-12 NCERTs for History, Geography, Polity, and Economy.",
            "**Phase 2: Standard Reference Books (Months 4-8):** Laxmikanth (Polity), Spectrum (History)."
        ],
        "top_teachers": [
            {"subject": "Indian Polity", "teacher": "Dr. Vikas Divyakirti", "channel": "Drishti IAS", "video_title": "Polity Concept Series", "youtube_link": "https://www.youtube.com/results?search_query=upsc+polity+drishti+ias"}
        ],
        "pyqs": [
            {"year": "2024", "title": "UPSC CSE Prelims General Studies (Paper-I & CSAT)", "link": "https://www.upsc.gov.in/", "format": "Official PDF"}
        ],
        "mocks": [
            {"question": "Which Schedule of the Indian Constitution contains the anti-defection law?", "options": ["8th Schedule", "9th Schedule", "10th Schedule", "11th Schedule"], "answer": "10th Schedule", "explanation": "The 10th Schedule was added by the 52nd Amendment Act in 1985."},
            {"question": "Which river is known as the 'Sorrow of Bengal'?", "options": ["Damodar River", "Hooghly River", "Kosi River", "Brahmaputra River"], "answer": "Damodar River", "explanation": "Damodar River was historically known as the Sorrow of Bengal due to ravaging floods."}
        ]
    },
    "NEET UG": {
        "full_name": "National Eligibility cum Entrance Test",
        "icon": "🩺",
        "category": "Medical Entrance",
        "description": "National entrance exam for admission into MBBS, BDS, and AYUSH courses.",
        "official_pyq_portal": "https://neet.nta.nic.in/document-category/archive/",
        "roadmap": [
            "**Phase 1: NCERT Biology & Chemistry (Months 1-6):** Line-by-line mastery of Class 11 & 12 NCERTs.",
            "**Phase 2: Physics Problem Solving (Months 7-9):** Daily numerical practice."
        ],
        "top_teachers": [
            {"subject": "Biology", "teacher": "Dr. Anand Mani", "channel": "Dr. Anand Mani", "video_title": "NCERT Biology One-Shot", "youtube_link": "https://www.youtube.com/results?search_query=neet+biology+one+shot"}
        ],
        "pyqs": [
            {"year": "2024", "title": "NEET UG Official Question Paper with Solutions", "link": "https://neet.nta.nic.in/", "format": "Official PDF"}
        ],
        "mocks": [
            {"question": "Which organelle is known as the powerhouse of the cell?", "options": ["Ribosome", "Mitochondria", "Golgi Apparatus", "Lysosome"], "answer": "Mitochondria", "explanation": "Mitochondria generate ATP."},
            {"question": "What is the pH of human blood under normal physiological conditions?", "options": ["7.35 - 7.45", "6.0 - 6.5", "8.0 - 8.5", "5.5 - 6.0"], "answer": "7.35 - 7.45", "explanation": "Human blood pH is strictly regulated between 7.35 and 7.45."}
        ]
    },
    "JEE Main & Advanced": {
        "full_name": "Joint Entrance Examination",
        "icon": "📐",
        "category": "Engineering Entrance",
        "description": "Gateway to IITs, NITs, IIITs, and premier engineering institutes.",
        "official_pyq_portal": "https://nta.ac.in/Downloads",
        "roadmap": [
            "**Phase 1: Foundation Building (Months 1-6):** Mechanics, Organic Chemistry, Calculus.",
            "**Phase 2: Advanced Problem Solving (Months 7-10):** Standard problem books."
        ],
        "top_teachers": [
            {"subject": "Mathematics", "teacher": "NV Sir", "channel": "Unacademy Atoms", "video_title": "BounceBack Series - Maths", "youtube_link": "https://www.youtube.com/results?search_query=bounceback+maths+nv+sir"}
        ],
        "pyqs": [
            {"year": "2024", "title": "JEE Main Session 1 & Session 2 All Papers", "link": "https://nta.ac.in/Downloads", "format": "Official PDF Archive"}
        ],
        "mocks": [
            {"question": "What is the derivative of sin(x^2) with respect to x?", "options": ["2x * cos(x^2)", "cos(x^2)", "-2x * cos(x^2)", "2 * sin(x)"], "answer": "2x * cos(x^2)", "explanation": "Applies the chain rule derivative: d/dx[sin(u)] = cos(u) * du/dx."}
        ]
    }
}

AGE_GROUPS = {
    "14 - 18 Years (School Students)": "Focus on balancing school/boards with foundational exam preparation.",
    "18 - 22 Years (College / Undergrads)": "Focus on parallel study schedules, college exams, and core conceptual preparation.",
    "22 - 28 Years (Working / Full-Time Prep)": "Focus on intense high-efficiency sessions, daily mock drills, and speed mastery.",
    "28+ Years (Senior Aspirants)": "Focus on high-yield selective studying, strategic time management, and mock analysis."
}


# ==========================================
# 4. SIDEBAR CONFIGURATIONS
# ==========================================
st.sidebar.title("⚙️ Personalization Hub")
st.sidebar.markdown("### 🎛️ Configure Your Profile")

exam_list = list(EXAMS_DATA.keys())

# 1. Standard Custom Exam Selection Dropdown
st.sidebar.markdown("**1. Select Target Exam:**")
selected_exam_name = st.sidebar.selectbox(
    "Select Target Exam",
    options=exam_list,
    index=0,
    format_func=lambda x: f"{EXAMS_DATA[x]['icon']} {x}",
    label_visibility="collapsed"
)

# 2. Age Group Selector
st.sidebar.markdown("---")
st.sidebar.markdown("**2. Aspirant Age Group / Category:**")
selected_age_group = st.sidebar.selectbox(
    "Age Group Selector",
    options=list(AGE_GROUPS.keys()),
    index=1,
    label_visibility="collapsed"
)

# 3. Daily Hours Capacity Slider
st.sidebar.markdown("---")
st.sidebar.markdown("**3. Daily Available Time (Hours):**")
daily_hours = st.sidebar.slider("Daily Hours Slider", 2, 14, 6, label_visibility="collapsed")

# 4. Months Remaining Slider
st.sidebar.markdown("**4. Months Remaining for Exam:**")
prep_months = st.sidebar.slider("Prep Months Slider", 1, 24, 6, label_visibility="collapsed")

# 5. Number of MCQs Selector
st.sidebar.markdown("---")
st.sidebar.markdown("**5. How many MCQs do you want?**")
num_mcqs = st.sidebar.slider("MCQ Count Slider", 1, 20, 5, label_visibility="collapsed")

st.sidebar.markdown("---")
st.sidebar.caption("🔑 **Optional Groq AI Key:**")
api_key = st.sidebar.text_input("Groq API Key", type="password", label_visibility="collapsed")

exam_info = EXAMS_DATA[selected_exam_name]


# ==========================================
# 5. MAIN DASHBOARD CONTENT & NAVIGATION
# ==========================================
# Header Banner
st.markdown(f"""
    <div class="glass-card">
        <span class="badge badge-primary">{exam_info['category']}</span>
        <span class="badge badge-purple">{selected_age_group}</span>
        <h1 style="margin-top: 10px; margin-bottom: 5px;">{exam_info['icon']} {selected_exam_name}</h1>
        <h4 style="color: #94a3b8; font-weight: 400; margin-bottom: 10px;">{exam_info['full_name']}</h4>
        <p style="color: #cbd5e1; font-size: 15px;">{exam_info['description']}</p>
    </div>
""", unsafe_allow_html=True)

# Tabs
tab_goal, tab_roadmap, tab_content, tab_pyqs, tab_mocks, tab_chat = st.tabs([
    "🎯 Custom Goal Planner",
    "🗺️ Roadmap & Strategy", 
    "📺 Curated Video Lectures", 
    "📄 PYQ Vault", 
    "⚡ Interactive Mock Test Engine",
    "💬 Ask Doubts AI"
])

# ------------------------------------------
# TAB 1: CUSTOM GOAL PLANNER
# ------------------------------------------
with tab_goal:
    st.subheader("🤖 Generate & Download Personal Goal Plan")
    st.markdown("Your study routine will be dynamically customized based on the parameters set in your sidebar.")
    
    col_info1, col_info2, col_info3 = st.columns(3)
    col_info1.metric("Selected Exam", selected_exam_name)
    col_info2.metric("Target Profile", selected_age_group.split(' ')[0] + " Yrs")
    col_info3.metric("Daily Schedule", f"{daily_hours} Hrs / day")

    if st.button("🚀 Generate Customized Goal Plan", use_container_width=True):
        if not api_key:
            plan_content = f"""Target Exam: {selected_exam_name} ({exam_info['full_name']})
Aspirant Profile: {selected_age_group}
Timeline: {prep_months} Months Out | Daily Schedule: {daily_hours} Hours/day

• Profile Strategy: {AGE_GROUPS[selected_age_group]}

• Morning Session ({int(daily_hours * 0.4)} hrs): Core Technical / High-Priority Conceptual Subject.
• Afternoon Session ({int(daily_hours * 0.3)} hrs): Secondary Subject Practice & Speed Building.
• Evening Session ({int(daily_hours * 0.2)} hrs): Previous Year Questions (PYQs) & Sectional Tests.
• Night Review ({round(daily_hours * 0.1, 1)} hrs): Daily Current Affairs, Formulas & Mistake Notebook Revision."""
        else:
            try:
                from groq import Groq
                client = Groq(api_key=api_key)
                prompt = f"Create a detailed daily study plan for {selected_exam_name} tailored for a {selected_age_group} student with {daily_hours} hours available daily and {prep_months} months remaining."
                response = client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role": "user", "content": prompt}]
                )
                plan_content = response.choices[0].message.content
            except Exception as e:
                st.error(f"Error calling Groq API: {str(e)}")
                plan_content = "Default Goal Plan Generated."

        st.session_state['generated_plan'] = plan_content
        st.session_state['plan_exam'] = selected_exam_name
        st.session_state['plan_age'] = selected_age_group
        st.session_state['plan_hours'] = daily_hours
        st.session_state['plan_months'] = prep_months

    if 'generated_plan' in st.session_state:
        st.markdown("---")
        st.markdown(f"### 📋 Personal Study Goal Plan for {st.session_state['plan_exam']}")
        st.text_area("Your Generated Plan", st.session_state['generated_plan'], height=220)

        st.markdown("### 📥 Download Plan File")
        col_doc, col_ppt = st.columns(2)

        docx_file = create_docx_plan(
            st.session_state['plan_exam'],
            st.session_state['plan_age'],
            st.session_state['plan_hours'],
            st.session_state['plan_months'],
            st.session_state['generated_plan']
        )
        col_doc.download_button(
            label="📄 Download Word Doc (.docx)",
            data=docx_file,
            file_name=f"{st.session_state['plan_exam']}_Goal_Plan.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )

        pptx_file = create_pptx_plan(
            st.session_state['plan_exam'],
            st.session_state['plan_age'],
            st.session_state['plan_hours'],
            st.session_state['plan_months'],
            st.session_state['generated_plan']
        )
        col_ppt.download_button(
            label="📊 Download PowerPoint (.pptx)",
            data=pptx_file,
            file_name=f"{st.session_state['plan_exam']}_Goal_Plan.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

        st.markdown("### 📲 Share Goal Plan")
        share_text = f"My Study Goal Plan for {st.session_state['plan_exam']}:\n\n{st.session_state['generated_plan']}"
        encoded_text = urllib.parse.quote(share_text)
        
        whatsapp_url = f"https://api.whatsapp.com/send?text={encoded_text}"
        email_url = f"mailto:?subject={urllib.parse.quote('My Study Goal Plan')}&body={encoded_text}"

        col_wa, col_em = st.columns(2)
        col_wa.markdown(f'<a href="{whatsapp_url}" target="_blank" class="share-btn-wa">💬 Share via WhatsApp</a>', unsafe_allow_html=True)
        col_em.markdown(f'<a href="{email_url}" target="_blank" class="share-btn-email">✉️ Share via Email</a>', unsafe_allow_html=True)

# ------------------------------------------
# TAB 2: ROADMAP
# ------------------------------------------
with tab_roadmap:
    st.subheader(f"📌 Preparation Strategy - {selected_exam_name}")
    for step in exam_info["roadmap"]:
        st.markdown(f"<div class='roadmap-step'>{step}</div>", unsafe_allow_html=True)

# ------------------------------------------
# TAB 3: VIDEO LECTURES
# ------------------------------------------
with tab_content:
    st.subheader(f"📺 Master Class Video Lectures - {selected_exam_name}")
    cols = st.columns(2)
    for idx, item in enumerate(exam_info["top_teachers"]):
        col = cols[idx % 2]
        with col:
            st.markdown(f"""
                <div class="glass-card">
                    <h4>📖 {item['subject']}</h4>
                    <p style="margin-bottom: 4px;">👨‍🏫 <b>Educator:</b> {item['teacher']}</p>
                    <p style="margin-bottom: 4px;">📢 <b>Channel:</b> {item['channel']}</p>
                    <p style="margin-bottom: 10px;">🎥 <b>Title:</b> {item['video_title']}</p>
                    <a href="{item['youtube_link']}" target="_blank" class="yt-btn">▶️ Watch Playlist on YouTube</a>
                </div>
            """, unsafe_allow_html=True)

# ------------------------------------------
# TAB 4: PYQs
# ------------------------------------------
with tab_pyqs:
    st.subheader(f"📄 Official PYQ Archive - {selected_exam_name}")
    st.markdown(f'<a href="{exam_info["official_pyq_portal"]}" target="_blank" class="official-btn">🌐 Open Official Exam Portal</a>', unsafe_allow_html=True)
    
    for pyq in exam_info["pyqs"]:
        st.markdown(f"""
            <div class="glass-card" style="display: flex; justify-content: space-between; align-items: center;">
                <div>
                    <h4>{pyq['title']} ({pyq['year']})</h4>
                    <span class="badge badge-success">{pyq['format']}</span>
                </div>
                <a href="{pyq['link']}" target="_blank" class="download-link-btn">⬇️ Download PDF</a>
            </div>
        """, unsafe_allow_html=True)

# ------------------------------------------
# TAB 5: MOCK TEST ENGINE (WITH DYNAMIC MCQ COUNT)
# ------------------------------------------
with tab_mocks:
    st.subheader(f"⚡ Practice Mock Test - {selected_exam_name}")
    st.info(f"Showing **{num_mcqs}** MCQ(s) as requested in sidebar settings.")

    base_questions = exam_info["mocks"]
    
    # Expand or pad questions dynamically to meet requested num_mcqs
    questions_to_show = []
    while len(questions_to_show) < num_mcqs:
        questions_to_show.extend(base_questions)
    questions_to_show = questions_to_show[:num_mcqs]

    with st.form("mock_test_form"):
        user_answers = {}
        for idx, q in enumerate(questions_to_show):
            st.markdown(f"**Q{idx+1}: {q['question']}**")
            user_answers[idx] = st.radio(
                f"Select option for Question {idx+1}", 
                q["options"], 
                key=f"q_{idx}",
                label_visibility="collapsed"
            )
            st.markdown("---")
            
        submitted = st.form_submit_button("📝 Submit Answers")

    if submitted:
        score = 0
        st.markdown("### 📊 Test Results & Analysis")
        for idx, q in enumerate(questions_to_show):
            ans = user_answers[idx]
            if ans == q["answer"]:
                score += 1
                st.success(f"**Q{idx+1}: Correct!** Choice: `{ans}`")
            else:
                st.error(f"**Q{idx+1}: Incorrect.** Your Choice: `{ans}` | Correct Answer: `{q['answer']}`")
            st.caption(f"💡 **Explanation:** {q['explanation']}")
            
        st.balloons()
        st.metric("Final Score", f"{score} / {num_mcqs}", f"{(score/num_mcqs)*100:.1f}%")

# ------------------------------------------
# TAB 6: DOUBT CHATBOX & VOICE ASSISTANT
# ------------------------------------------
with tab_chat:
    st.subheader("💬 Ask Doubts & AI Voice Assistant")
    st.markdown("Clear your doubts instantly using text or voice inputs!")

    # Initialize chat history
    if "chat_messages" not in st.session_state:
        st.session_state.chat_messages = [
            {"role": "assistant", "content": f"Hello! I am your AI Tutor for **{selected_exam_name}**. Ask me any doubt or concept you want explained!"}
        ]

    # Render previous messages
    for msg in st.session_state.chat_messages:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

    st.markdown("---")
    
    # Audio input option
    st.markdown("#### 🎙️ Voice Input (Record Your Doubt)")
    audio_value = st.audio_input("Record audio doubt")

    # Text input option
    user_doubt = st.chat_input(f"Type your doubt regarding {selected_exam_name}...")

    prompt_to_process = None

    if audio_value:
        prompt_to_process = f"[Voice Input Received] Can you explain the general key concepts for {selected_exam_name}?"
    elif user_doubt:
        prompt_to_process = user_doubt

    if prompt_to_process:
        # Append User Message
        st.session_state.chat_messages.append({"role": "user", "content": prompt_to_process})
        with st.chat_message("user"):
            st.write(prompt_to_process)

        # Generate AI Response
        with st.chat_message("assistant"):
            with st.spinner("Analyzing doubt & preparing response..."):
                if api_key:
                    try:
                        from groq import Groq
                        client = Groq(api_key=api_key)
                        response = client.chat.completions.create(
                            model="llama-3.3-70b-versatile",
                            messages=[
                                {"role": "system", "content": f"You are an expert tutor helping a student prepare for {selected_exam_name}. Keep explanations concise, accurate, and easy to understand."},
                                *[{"role": m["role"], "content": m["content"]} for m in st.session_state.chat_messages]
                            ]
                        )
                        reply = response.choices[0].message.content
                    except Exception as e:
                        reply = f"Error generating answer: {str(e)}"
                else:
                    reply = f"Here is a quick guidance on your query regarding **{selected_exam_name}**:\n\nFocus on core syllabus topics, revise previous year questions, and ensure daily speed tests. *(Add a Groq API Key in the sidebar for full conversational AI answers!)*"

                st.write(reply)
                
                # Audio response synthesis (Voice Output)
                if HAS_GTTS:
                    try:
                        audio_fp = text_to_speech(reply[:250])  # Convert first 250 chars to speech
                        if audio_fp:
                            st.audio(audio_fp, format="audio/mp3")
                    except Exception as audio_err:
                        st.caption("🔊 Voice response generation skipped.")

        st.session_state.chat_messages.append({"role": "assistant", "content": reply})
